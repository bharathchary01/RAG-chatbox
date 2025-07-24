import os, csv, docx, fitz
from pptx import Presentation
from mcp.message_protocol import Message

class IngestionAgent:
    def ingest(self, documents, trace_id):
        chunks = []
        for path in documents:
            ext = os.path.splitext(path)[1].lower()
            if ext == ".txt":
                with open(path) as f:
                    chunks += f.read().splitlines()
            elif ext == ".pdf":
                doc = fitz.open(path)
                chunks += [pg.get_text() for pg in doc]
            elif ext == ".docx":
                chunks += [p.text for p in docx.Document(path).paragraphs if p.text.strip()]
            elif ext == ".csv":
                with open(path) as f:
                    chunks += [", ".join(r) for r in csv.reader(f)]
            elif ext == ".pptx":
                prs = Presentation(path)
                for slide in prs.slides:
                    chunks.append(" ".join(shape.text for shape in slide.shapes if hasattr(shape, "text")))
        return Message("IngestionAgent", "RetrievalAgent", "DOCUMENT_PARSED", trace_id, {"chunks": chunks})
